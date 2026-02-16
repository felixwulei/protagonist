"""Local macOS tools — executed on the user's machine via the menu bar app."""
from __future__ import annotations

import os
import json
import asyncio
import base64
import tempfile
import subprocess


# --------------- Tool Registry ---------------

TOOL_NAMES = [
    "read_emails", "search_emails", "send_email",
    "capture_screen", "check_wechat",
    "run_claude_code",
    "find_files", "read_file",
    "open_app", "quit_app", "run_command",
    "get_calendar_events", "create_reminder",
    "music_play", "music_pause", "music_next", "music_previous",
    "music_now_playing", "music_search_play",
    "create_document",
]

# Tool definitions for OpenAI function calling (sent to agent when local tools are available)
TOOL_DEFINITIONS = [
    {"type": "function", "function": {"name": "read_emails", "description": "Read recent emails from Mail.app", "parameters": {"type": "object", "properties": {"count": {"type": "integer", "description": "How many emails", "default": 5}, "unread_only": {"type": "boolean", "description": "Unread only", "default": False}}}}},
    {"type": "function", "function": {"name": "search_emails", "description": "Search emails by keyword", "parameters": {"type": "object", "properties": {"query": {"type": "string", "description": "Search keyword"}, "count": {"type": "integer", "default": 5}}, "required": ["query"]}}},
    {"type": "function", "function": {"name": "send_email", "description": "Send an email via Mail.app", "parameters": {"type": "object", "properties": {"to": {"type": "string"}, "subject": {"type": "string"}, "body": {"type": "string"}}, "required": ["to", "subject", "body"]}}},
    {"type": "function", "function": {"name": "capture_screen", "description": "Screenshot the screen or a specific app and analyze what's visible", "parameters": {"type": "object", "properties": {"app_name": {"type": "string", "description": "App to capture (omit for full screen)"}, "question": {"type": "string", "description": "What to look for", "default": "Describe what's on screen"}}}}},
    {"type": "function", "function": {"name": "check_wechat", "description": "Screenshot WeChat and check for unread messages", "parameters": {"type": "object", "properties": {"question": {"type": "string", "default": "Any unread messages? Who sent them?"}}}}},
    {"type": "function", "function": {"name": "run_claude_code", "description": "Run Claude Code CLI for a programming task", "parameters": {"type": "object", "properties": {"task": {"type": "string", "description": "The programming task"}, "working_dir": {"type": "string", "description": "Working directory"}}, "required": ["task"]}}},
    {"type": "function", "function": {"name": "find_files", "description": "Search for files using Spotlight", "parameters": {"type": "object", "properties": {"query": {"type": "string"}, "file_type": {"type": "string", "description": "e.g. pdf, docx, py"}}, "required": ["query"]}}},
    {"type": "function", "function": {"name": "read_file", "description": "Read a file's contents", "parameters": {"type": "object", "properties": {"path": {"type": "string"}}, "required": ["path"]}}},
    {"type": "function", "function": {"name": "open_app", "description": "Open an application", "parameters": {"type": "object", "properties": {"app_name": {"type": "string"}}, "required": ["app_name"]}}},
    {"type": "function", "function": {"name": "quit_app", "description": "Quit an application", "parameters": {"type": "object", "properties": {"app_name": {"type": "string"}}, "required": ["app_name"]}}},
    {"type": "function", "function": {"name": "run_command", "description": "Run a shell command", "parameters": {"type": "object", "properties": {"command": {"type": "string"}, "timeout": {"type": "integer", "default": 30}}, "required": ["command"]}}},
    {"type": "function", "function": {"name": "get_calendar_events", "description": "Get upcoming calendar events", "parameters": {"type": "object", "properties": {"days": {"type": "integer", "default": 1}}}}},
    {"type": "function", "function": {"name": "create_reminder", "description": "Create a reminder", "parameters": {"type": "object", "properties": {"title": {"type": "string"}, "due_date": {"type": "string", "description": "YYYY-MM-DD HH:MM"}}, "required": ["title"]}}},
    {"type": "function", "function": {"name": "music_play", "description": "Resume or start playing music (Apple Music or Spotify)", "parameters": {"type": "object", "properties": {}}}},
    {"type": "function", "function": {"name": "music_pause", "description": "Pause the currently playing music", "parameters": {"type": "object", "properties": {}}}},
    {"type": "function", "function": {"name": "music_next", "description": "Skip to the next track", "parameters": {"type": "object", "properties": {}}}},
    {"type": "function", "function": {"name": "music_previous", "description": "Go back to the previous track", "parameters": {"type": "object", "properties": {}}}},
    {"type": "function", "function": {"name": "music_now_playing", "description": "Get info about the currently playing track", "parameters": {"type": "object", "properties": {}}}},
    {"type": "function", "function": {"name": "music_search_play", "description": "Search for a song, artist or album and play it", "parameters": {"type": "object", "properties": {"query": {"type": "string", "description": "Song name, artist, or album to search and play"}}, "required": ["query"]}}},
    {"type": "function", "function": {"name": "create_document", "description": "Create a document file (Word, PowerPoint slides, or PDF) from markdown content and send it to the user", "parameters": {"type": "object", "properties": {"type": {"type": "string", "enum": ["word", "slides", "pdf"], "description": "Document type: word (.docx), slides (.pptx), or pdf (.pdf)"}, "title": {"type": "string", "description": "Document title / filename (without extension)"}, "content": {"type": "string", "description": "Document content in markdown format. Use # for headings, ## for subheadings, - for bullet points, regular text for paragraphs. For slides: each # starts a new slide."}}, "required": ["type", "title", "content"]}}},
]


async def execute_tool(name: str, args: dict) -> str:
    """Execute a local tool and return text result."""
    print(f"  [tool] {name}({json.dumps(args, ensure_ascii=False)})")
    try:
        fn = _TOOL_MAP.get(name)
        if not fn:
            return f"Unknown tool: {name}"
        result = await fn(**args)
        if len(result) > 4000:
            result = result[:4000] + "\n...(truncated)"
        return result
    except Exception as e:
        return f"Tool error: {e}"


# --------------- Helpers ---------------

async def _applescript(script: str) -> str:
    proc = await asyncio.create_subprocess_exec(
        "osascript", "-e", script,
        stdout=subprocess.PIPE, stderr=subprocess.PIPE,
    )
    stdout, stderr = await proc.communicate()
    if proc.returncode != 0:
        return f"AppleScript error: {stderr.decode().strip()}"
    return stdout.decode().strip()


async def _shell(cmd: str, timeout: int = 30) -> str:
    try:
        proc = await asyncio.create_subprocess_shell(
            cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
        )
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
        output = stdout.decode() + stderr.decode()
        return output.strip() if output.strip() else "(no output)"
    except asyncio.TimeoutError:
        proc.kill()
        return "(timed out)"
    except Exception as e:
        return f"Error: {e}"


async def _screenshot_analyze(question: str = "Describe what's on screen", app_name: str = None) -> str:
    """Screenshot + GPT-4V analysis."""
    from openai import AsyncOpenAI
    client = AsyncOpenAI(api_key=os.getenv("OPENAI_API_KEY"))
    model = os.getenv("LLM_MODEL", "gpt-4o-mini")

    tmp = tempfile.mktemp(suffix=".png")
    try:
        if app_name:
            await _applescript(f'tell application "{app_name}" to activate')
            await asyncio.sleep(0.8)
        await _shell(f'screencapture -x "{tmp}"', timeout=5)
        if not os.path.exists(tmp):
            return "Screenshot failed"
        with open(tmp, "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
        response = await client.chat.completions.create(
            model=model, max_tokens=500,
            messages=[{"role": "user", "content": [
                {"type": "text", "text": question},
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}", "detail": "low"}},
            ]}],
        )
        return response.choices[0].message.content or "Could not analyze"
    finally:
        if os.path.exists(tmp):
            os.unlink(tmp)


# --------------- Email App Detection ---------------

# Microsoft Graph API — for New Outlook which doesn't support AppleScript
GRAPH_CLIENT_ID = "d3590ed6-52b3-4102-aeff-aad2292ab01c"  # Microsoft Office public client
GRAPH_SCOPES = ["Mail.Read", "Mail.Send", "User.Read", "offline_access"]
_GRAPH_CONFIG_PATH = os.path.expanduser("~/.protagonist/config.json")


def _load_graph_tokens() -> dict:
    """Load Graph API tokens from config."""
    try:
        with open(_GRAPH_CONFIG_PATH, "r") as f:
            return json.load(f).get("graph_tokens", {})
    except Exception:
        return {}


def _save_graph_tokens(tokens: dict):
    """Save Graph API tokens to config."""
    try:
        with open(_GRAPH_CONFIG_PATH, "r") as f:
            cfg = json.load(f)
    except Exception:
        cfg = {}
    cfg["graph_tokens"] = tokens
    with open(_GRAPH_CONFIG_PATH, "w") as f:
        json.dump(cfg, f, indent=2, ensure_ascii=False)


async def _graph_refresh_token() -> str | None:
    """Refresh the access token using the refresh token. Returns new access_token or None."""
    tokens = _load_graph_tokens()
    refresh = tokens.get("refresh_token")
    if not refresh:
        return None
    try:
        import urllib.request
        import urllib.parse
        data = urllib.parse.urlencode({
            "client_id": GRAPH_CLIENT_ID,
            "grant_type": "refresh_token",
            "refresh_token": refresh,
            "scope": " ".join(GRAPH_SCOPES),
        }).encode()
        req = urllib.request.Request(
            "https://login.microsoftonline.com/common/oauth2/v2.0/token",
            data=data,
            headers={"Content-Type": "application/x-www-form-urlencoded"},
        )
        resp = urllib.request.urlopen(req, timeout=15)
        result = json.loads(resp.read().decode())
        tokens["access_token"] = result["access_token"]
        if "refresh_token" in result:
            tokens["refresh_token"] = result["refresh_token"]
        _save_graph_tokens(tokens)
        return result["access_token"]
    except Exception as e:
        print(f"[graph] Token refresh failed: {e}")
        return None


async def _graph_get_token() -> str | None:
    """Get a valid Graph API access token (refresh if needed)."""
    tokens = _load_graph_tokens()
    access = tokens.get("access_token")
    if not access:
        return None
    # Try using the token; if it fails we'll refresh
    return access


async def _graph_api(endpoint: str, method: str = "GET", body: dict = None) -> dict | str:
    """Call Microsoft Graph API. Returns parsed JSON or error string."""
    token = await _graph_get_token()
    if not token:
        return "EMAIL_AUTH_NEEDED"

    import urllib.request
    import urllib.error
    url = f"https://graph.microsoft.com/v1.0{endpoint}"
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json",
    }

    try:
        if body:
            data = json.dumps(body).encode()
        else:
            data = None
        req = urllib.request.Request(url, data=data, headers=headers, method=method)
        resp = urllib.request.urlopen(req, timeout=20)
        return json.loads(resp.read().decode())
    except urllib.error.HTTPError as e:
        if e.code == 401:
            # Token expired, try refresh
            new_token = await _graph_refresh_token()
            if new_token:
                headers["Authorization"] = f"Bearer {new_token}"
                req = urllib.request.Request(url, data=data if body else None, headers=headers, method=method)
                try:
                    resp = urllib.request.urlopen(req, timeout=20)
                    return json.loads(resp.read().decode())
                except Exception as e2:
                    return f"Graph API error after refresh: {e2}"
            return "EMAIL_AUTH_NEEDED"
        return f"Graph API error: {e.code} {e.read().decode()[:200]}"
    except Exception as e:
        return f"Graph API error: {e}"


async def _graph_device_code_auth() -> str:
    """Start device code flow for Graph API auth. Returns instructions for the user."""
    import urllib.request
    import urllib.parse
    data = urllib.parse.urlencode({
        "client_id": GRAPH_CLIENT_ID,
        "scope": " ".join(GRAPH_SCOPES),
    }).encode()
    req = urllib.request.Request(
        "https://login.microsoftonline.com/common/oauth2/v2.0/devicecode",
        data=data,
        headers={"Content-Type": "application/x-www-form-urlencoded"},
    )
    resp = urllib.request.urlopen(req, timeout=15)
    result = json.loads(resp.read().decode())

    user_code = result["user_code"]
    device_code = result["device_code"]
    verification_uri = result.get("verification_uri", "https://microsoft.com/devicelogin")
    expires_in = result.get("expires_in", 900)
    interval = result.get("interval", 5)

    # Poll for token in background
    async def _poll():
        import time
        deadline = time.time() + expires_in
        while time.time() < deadline:
            await asyncio.sleep(interval)
            try:
                poll_data = urllib.parse.urlencode({
                    "client_id": GRAPH_CLIENT_ID,
                    "grant_type": "urn:ietf:params:oauth:grant-type:device_code",
                    "device_code": device_code,
                }).encode()
                poll_req = urllib.request.Request(
                    "https://login.microsoftonline.com/common/oauth2/v2.0/token",
                    data=poll_data,
                    headers={"Content-Type": "application/x-www-form-urlencoded"},
                )
                poll_resp = urllib.request.urlopen(poll_req, timeout=15)
                tokens = json.loads(poll_resp.read().decode())
                _save_graph_tokens({
                    "access_token": tokens["access_token"],
                    "refresh_token": tokens.get("refresh_token", ""),
                })
                print(f"[graph] Authorization successful!")
                return
            except urllib.error.HTTPError as e:
                body = e.read().decode()
                if "authorization_pending" in body:
                    continue
                elif "expired_token" in body:
                    print("[graph] Device code expired")
                    return
                else:
                    continue
            except Exception:
                continue

    asyncio.create_task(_poll())
    return (
        f"The user's email requires one-time authorization. "
        f"Tell them to open {verification_uri} and enter code: {user_code} — "
        f"then try reading email again in a minute. After they authorize, it'll work permanently."
    )


def _email_app() -> str:
    """Detect email app: graph tokens > config > Outlook AppleScript test > mail."""
    # 1. Check if Graph API tokens exist (New Outlook users)
    tokens = _load_graph_tokens()
    if tokens.get("access_token") or tokens.get("refresh_token"):
        return "graph"

    # 2. Check config preference
    try:
        with open(_GRAPH_CONFIG_PATH, "r") as f:
            pref = json.load(f).get("email_app", "")
        if pref in ("graph", "outlook", "mail"):
            return pref
    except Exception:
        pass

    # 3. Check if Outlook is installed
    if os.path.exists("/Applications/Microsoft Outlook.app"):
        return "outlook"
    return "mail"


async def _graph_read_emails(count: int = 5, unread_only: bool = False) -> str:
    """Read emails via Microsoft Graph API."""
    filter_str = "&$filter=isRead eq false" if unread_only else ""
    result = await _graph_api(
        f"/me/messages?$top={count}&$orderby=receivedDateTime desc"
        f"&$select=subject,from,receivedDateTime,isRead{filter_str}"
    )
    if isinstance(result, str):
        return result

    messages = result.get("value", [])
    if not messages:
        return "No emails found"

    lines = []
    for m in messages:
        unread = "[UNREAD] " if not m.get("isRead") else ""
        sender = m.get("from", {}).get("emailAddress", {})
        name = sender.get("name", "")
        addr = sender.get("address", "")
        subj = m.get("subject", "(no subject)")
        dt = m.get("receivedDateTime", "")[:16].replace("T", " ")
        lines.append(f"{unread}From: {name} <{addr}> | {subj} | {dt}")
    return "\n".join(lines)


async def _graph_search_emails(query: str, count: int = 5) -> str:
    """Search emails via Microsoft Graph API."""
    import urllib.parse
    q = urllib.parse.quote(query)
    result = await _graph_api(
        f"/me/messages?$top={count}&$search=\"{q}\""
        f"&$select=subject,from,receivedDateTime,isRead"
    )
    if isinstance(result, str):
        return result

    messages = result.get("value", [])
    if not messages:
        return "No matching emails"

    lines = []
    for m in messages:
        sender = m.get("from", {}).get("emailAddress", {})
        name = sender.get("name", "")
        addr = sender.get("address", "")
        subj = m.get("subject", "(no subject)")
        dt = m.get("receivedDateTime", "")[:16].replace("T", " ")
        lines.append(f"From: {name} <{addr}> | {subj} | {dt}")
    return "\n".join(lines)


async def _graph_send_email(to: str, subject: str, body: str) -> str:
    """Send email via Microsoft Graph API."""
    result = await _graph_api("/me/sendMail", method="POST", body={
        "message": {
            "subject": subject,
            "body": {"contentType": "Text", "content": body},
            "toRecipients": [{"emailAddress": {"address": to}}],
        }
    })
    if isinstance(result, str):
        return result
    return f"Email sent to {to}"


# --------------- Tool Implementations ---------------

async def _read_emails(count: int = 5, unread_only: bool = False) -> str:
    app = _email_app()
    if app == "graph":
        result = await _graph_read_emails(count, unread_only)
        if result == "EMAIL_AUTH_NEEDED":
            auth_result = await _graph_device_code_auth()
            return auth_result
        return result
    elif app == "outlook":
        f = "whose is read is false" if unread_only else ""
        result = await _applescript(f'''
        tell application "Microsoft Outlook"
            set output to ""
            set msgs to messages of inbox {f}
            set n to {count}
            if (count of msgs) < n then set n to (count of msgs)
            repeat with i from 1 to n
                set m to item i of msgs
                set isRead to is read of m
                set mk to ""
                if not isRead then set mk to "[UNREAD] "
                set senderName to name of sender of m
                set senderAddr to address of sender of m
                set output to output & mk & "From: " & senderName & " <" & senderAddr & "> | " & subject of m & " | " & (time received of m as string) & linefeed
            end repeat
            return output
        end tell''') or "No emails found"
        # If Outlook returned nothing, it might be New Outlook — try Graph
        if result == "No emails found" or result.strip() == "":
            graph_result = await _graph_read_emails(count, unread_only)
            if graph_result == "EMAIL_AUTH_NEEDED":
                auth_result = await _graph_device_code_auth()
                return auth_result
            if graph_result and "error" not in graph_result.lower():
                # Auto-switch to graph for future calls
                try:
                    with open(_GRAPH_CONFIG_PATH, "r") as cfg_f:
                        cfg = json.load(cfg_f)
                    cfg["email_app"] = "graph"
                    with open(_GRAPH_CONFIG_PATH, "w") as cfg_f:
                        json.dump(cfg, cfg_f, indent=2, ensure_ascii=False)
                except Exception:
                    pass
                return graph_result
        return result
    else:
        f = "whose read status is false" if unread_only else ""
        return await _applescript(f'''
        tell application "Mail"
            set output to ""
            set msgs to (messages of inbox {f})
            set n to {count}
            if (count of msgs) < n then set n to (count of msgs)
            repeat with i from 1 to n
                set m to item i of msgs
                set isRead to read status of m
                set mk to ""
                if not isRead then set mk to "[UNREAD] "
                set output to output & mk & "From: " & sender of m & " | " & subject of m & " | " & (date received of m as string) & linefeed
            end repeat
            return output
        end tell''') or "No emails found"


async def _search_emails(query: str, count: int = 5) -> str:
    q = query.replace('"', '\\"')
    app = _email_app()
    if app == "graph":
        result = await _graph_search_emails(query, count)
        if result == "EMAIL_AUTH_NEEDED":
            auth_result = await _graph_device_code_auth()
            return auth_result
        return result
    elif app == "outlook":
        return await _applescript(f'''
        tell application "Microsoft Outlook"
            set output to ""
            set results to messages of inbox whose subject contains "{q}" or (address of sender contains "{q}") or (name of sender contains "{q}")
            set n to {count}
            if (count of results) < n then set n to (count of results)
            repeat with i from 1 to n
                set m to item i of results
                set senderName to name of sender of m
                set senderAddr to address of sender of m
                set output to output & "From: " & senderName & " <" & senderAddr & "> | " & subject of m & " | " & (time received of m as string) & linefeed
            end repeat
            if output is "" then return "No matching emails"
            return output
        end tell''')
    else:
        return await _applescript(f'''
        tell application "Mail"
            set output to ""
            set results to (messages of inbox whose subject contains "{q}" or sender contains "{q}")
            set n to {count}
            if (count of results) < n then set n to (count of results)
            repeat with i from 1 to n
                set m to item i of results
                set output to output & "From: " & sender of m & " | " & subject of m & " | " & (date received of m as string) & linefeed
            end repeat
            if output is "" then return "No matching emails"
            return output
        end tell''')


async def _send_email(to: str, subject: str, body: str) -> str:
    s = subject.replace('"', '\\"')
    b = body.replace('"', '\\"')
    app = _email_app()
    if app == "graph":
        result = await _graph_send_email(to, subject, body)
        if result == "EMAIL_AUTH_NEEDED":
            auth_result = await _graph_device_code_auth()
            return auth_result
        return result
    elif app == "outlook":
        result = await _applescript(f'''
        tell application "Microsoft Outlook"
            set newMsg to make new outgoing message with properties {{subject:"{s}", content:"{b}"}}
            make new to recipient of newMsg with properties {{email address:{{address:"{to}"}}}}
            send newMsg
        end tell''')
    else:
        result = await _applescript(f'''
        tell application "Mail"
            set newMsg to make new outgoing message with properties {{subject:"{s}", content:"{b}", visible:true}}
            tell newMsg to make new to recipient at end of to recipients with properties {{address:"{to}"}}
            send newMsg
        end tell''')
    return f"Email sent to {to}" if "error" not in result.lower() else f"Failed: {result}"


async def _capture_screen(app_name: str = None, question: str = "Describe what's on screen") -> str:
    return await _screenshot_analyze(question=question, app_name=app_name)


async def _check_wechat(question: str = "Any unread messages? Who sent them?") -> str:
    check = await _applescript('tell application "System Events" to (name of processes) contains "WeChat"')
    if "true" not in check.lower():
        return "WeChat is not running"
    return await _screenshot_analyze(f"This is WeChat. {question}", app_name="WeChat")


async def _run_claude_code(task: str, working_dir: str = None) -> str:
    cmd = f'claude -p {json.dumps(task, ensure_ascii=False)} --output-format text'
    if working_dir:
        cmd = f'cd {json.dumps(working_dir)} && {cmd}'
    print(f"  [claude-code] {task[:80]}...")
    result = await _shell(cmd, timeout=120)
    print(f"  [claude-code] Done ({len(result)} chars)")
    return result


async def _find_files(query: str, file_type: str = None) -> str:
    cmd = f'mdfind "{query}" | grep -i "\\.{file_type}$" | head -15' if file_type else f'mdfind "{query}" | head -15'
    return await _shell(cmd, timeout=10) or f"No files found for '{query}'"


async def _read_file(path: str) -> str:
    expanded = os.path.expanduser(path)
    if not os.path.exists(expanded):
        return f"File not found: {path}"
    ext = os.path.splitext(expanded)[1].lower()
    if ext in (".pdf", ".docx", ".doc", ".rtf"):
        return await _shell(f'textutil -stdout -convert txt "{expanded}" 2>/dev/null || strings "{expanded}" | head -100')
    try:
        with open(expanded, "r", errors="replace") as f:
            return f.read(8000)
    except Exception as e:
        return f"Read failed: {e}"


async def _open_app(app_name: str) -> str:
    await _applescript(f'tell application "{app_name}" to activate')
    return f"Opened {app_name}"


async def _quit_app(app_name: str) -> str:
    result = await _applescript(f'tell application "{app_name}" to quit')
    return f"Quit {app_name}" if "error" not in result.lower() else f"Could not quit {app_name}"


async def _run_command(command: str, timeout: int = 30) -> str:
    return await _shell(command, timeout=timeout)


async def _get_calendar_events(days: int = 1) -> str:
    return await _applescript(f'''
    set today to current date
    set endDate to today + ({days} * days)
    set output to ""
    tell application "Calendar"
        repeat with cal in calendars
            set evts to (every event of cal whose start date >= today and start date <= endDate)
            repeat with evt in evts
                set output to output & summary of evt & " | " & (start date of evt as string) & " - " & (end date of evt as string) & linefeed
            end repeat
        end repeat
    end tell
    if output is "" then return "No events in the next {days} day(s)"
    return output''') or f"No events in the next {days} day(s)"


async def _create_reminder(title: str, due_date: str = None) -> str:
    t = title.replace('"', '\\"')
    if due_date:
        # Convert YYYY-MM-DD HH:MM to AppleScript-friendly format
        # AppleScript needs the system locale format; safest is to parse and reconstruct
        try:
            from datetime import datetime as _dt
            # Accept various formats
            for fmt in ("%Y-%m-%d %H:%M", "%Y-%m-%d", "%Y/%m/%d %H:%M", "%Y/%m/%d"):
                try:
                    parsed = _dt.strptime(due_date, fmt)
                    break
                except ValueError:
                    continue
            else:
                parsed = None

            if parsed:
                # Use AppleScript's date parsing with explicit month name (English, universally accepted)
                months = ["January", "February", "March", "April", "May", "June",
                          "July", "August", "September", "October", "November", "December"]
                as_date = f'{months[parsed.month - 1]} {parsed.day}, {parsed.year} {parsed.strftime("%I:%M:%S %p")}'
                script = f'tell application "Reminders" to make new reminder with properties {{name:"{t}", due date:date "{as_date}"}}'
            else:
                # Fallback: let AppleScript try the raw string
                script = f'tell application "Reminders" to make new reminder with properties {{name:"{t}", due date:date "{due_date}"}}'
        except Exception:
            script = f'tell application "Reminders" to make new reminder with properties {{name:"{t}"}}'
    else:
        script = f'tell application "Reminders" to make new reminder with properties {{name:"{t}"}}'
    result = await _applescript(script)
    return f"Reminder created: {title}" + (f" (due: {due_date})" if due_date else "") if "error" not in result.lower() else f"Failed: {result}"


def _music_app() -> str:
    """Detect which music app is running: Music (Apple Music) or Spotify."""
    # Check Spotify first since it's more common for global users
    import subprocess as sp
    result = sp.run(
        ["osascript", "-e", 'tell application "System Events" to (name of processes) contains "Spotify"'],
        capture_output=True, text=True, timeout=3,
    )
    if "true" in result.stdout.lower():
        return "Spotify"
    return "Music"


async def _music_play() -> str:
    app = _music_app()
    await _applescript(f'tell application "{app}" to activate')
    await asyncio.sleep(0.3)
    # If nothing is queued, play the user's library
    result = await _applescript(f'''
    tell application "{app}"
        if player state is not playing then
            try
                play playlist "Library"
            on error
                play
            end try
        else
            play
        end if
    end tell''')
    await asyncio.sleep(0.5)
    return await _music_now_playing()


async def _music_pause() -> str:
    app = _music_app()
    await _applescript(f'tell application "{app}" to pause')
    return f"Paused {app}"


async def _music_next() -> str:
    app = _music_app()
    await _applescript(f'tell application "{app}" to next track')
    # Get new track info
    await asyncio.sleep(0.5)
    return await _music_now_playing()


async def _music_previous() -> str:
    app = _music_app()
    await _applescript(f'tell application "{app}" to previous track')
    await asyncio.sleep(0.5)
    return await _music_now_playing()


async def _music_now_playing() -> str:
    app = _music_app()
    if app == "Spotify":
        result = await _applescript('''
        tell application "Spotify"
            if player state is playing then
                set t to name of current track
                set a to artist of current track
                set al to album of current track
                return "Playing: " & t & " by " & a & " (Album: " & al & ")"
            else
                return "Not playing"
            end if
        end tell''')
    else:
        result = await _applescript('''
        tell application "Music"
            if player state is playing then
                set t to name of current track
                set a to artist of current track
                set al to album of current track
                return "Playing: " & t & " by " & a & " (Album: " & al & ")"
            else
                return "Not playing"
            end if
        end tell''')
    return result or "No music app is playing"


async def _music_search_play(query: str) -> str:
    app = _music_app()
    q = query.replace('"', '\\"')
    if app == "Spotify":
        uri = f"spotify:search:{q}"
        await _shell(f'open "{uri}"', timeout=5)
        await asyncio.sleep(1.5)
        await _applescript('''
        tell application "Spotify" to activate
        delay 0.5
        tell application "System Events"
            keystroke return
        end tell''')
        await asyncio.sleep(1)
        return await _music_now_playing()
    else:
        import urllib.parse

        # Step 1: look up canonical track name via iTunes API
        encoded = urllib.parse.quote_plus(query)
        search_names = [q]  # always search the original query
        api_name, api_artist, api_url = "", "", ""
        # Search both CN and US stores to get both Chinese and English names
        for country in ["CN", "US"]:
            api_result = await _shell(
                f'curl -s "https://itunes.apple.com/search?term={encoded}&media=music&limit=1&country={country}"',
                timeout=10,
            )
            try:
                data = json.loads(api_result)
                if data.get("resultCount", 0) > 0:
                    track = data["results"][0]
                    name = track.get("trackName", "")
                    artist = track.get("artistName", "")
                    url = track.get("trackViewUrl", "")
                    if not api_name:
                        api_name = name
                    if not api_artist:
                        api_artist = artist
                    if not api_url:
                        api_url = url
                    for n in [name, artist]:
                        clean = n.replace('"', '\\"')
                        if clean and clean not in search_names:
                            search_names.append(clean)
            except (json.JSONDecodeError, KeyError):
                pass

        # Step 2: search ALL playlists with each name variant separately
        result = "NOT_FOUND"
        for search_term in search_names:
            st = search_term.replace('"', '\\"')
            print(f"  [music] Searching playlists for: {st}")
            r = await _applescript(
                f'tell application "Music"\n'
                f'activate\n'
                f'set allPlaylists to every playlist\n'
                f'repeat with p in allPlaylists\n'
                f'try\n'
                f'set results to (every track of p whose name contains "{st}")\n'
                f'if (count of results) > 0 then\n'
                f'set t to item 1 of results\n'
                f'play t\n'
                f'return "Playing: " & name of t & " by " & artist of t\n'
                f'end if\n'
                f'end try\n'
                f'try\n'
                f'set results to (every track of p whose artist contains "{st}")\n'
                f'if (count of results) > 0 then\n'
                f'set t to item 1 of results\n'
                f'play t\n'
                f'return "Playing: " & name of t & " by " & artist of t\n'
                f'end if\n'
                f'end try\n'
                f'end repeat\n'
                f'return "NOT_FOUND"\n'
                f'end tell'
            )
            print(f"  [music] Result: {r[:100]}")
            if r.startswith("Playing:"):
                result = r
                break
        if "NOT_FOUND" not in result:
            return result

        # Step 3: not in local playlists — search in Apple Music catalog
        search_term = api_name or query
        st = search_term.replace('"', '\\"')
        print(f"  [music] Not in playlists, searching Apple Music catalog for: {st}")
        r = await _applescript(
            f'tell application "Music"\n'
            f'activate\n'
            f'end tell\n'
            f'delay 0.3\n'
            f'tell application "System Events"\n'
            f'tell process "Music"\n'
            f'-- Open search via Cmd+F or click search field\n'
            f'keystroke "f" using {{command down, option down}}\n'
            f'delay 0.5\n'
            f'keystroke "{st}"\n'
            f'delay 0.3\n'
            f'keystroke return\n'
            f'end tell\n'
            f'end tell'
        )
        await asyncio.sleep(2)
        # Try to play the first result via Enter
        await _applescript(
            'tell application "System Events"\n'
            'tell process "Music"\n'
            'keystroke return\n'
            'end tell\n'
            'end tell'
        )
        await asyncio.sleep(1)
        now = await _music_now_playing()
        if "Playing:" in now:
            return now
        if api_name:
            return f"Searched for '{api_name}' in Apple Music app. Check the app to play."
        return f"Searched for '{query}' in Apple Music app. Check the app to play."


# --------------- Document Creation ---------------

_DOC_DIR = "/tmp/protagonist_docs"


def _ensure_doc_dir():
    os.makedirs(_DOC_DIR, exist_ok=True)


def _parse_markdown_blocks(content: str) -> list[dict]:
    """Parse markdown into blocks: heading1, heading2, heading3, bullet, paragraph."""
    blocks = []
    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("### "):
            blocks.append({"type": "heading3", "text": stripped[4:]})
        elif stripped.startswith("## "):
            blocks.append({"type": "heading2", "text": stripped[3:]})
        elif stripped.startswith("# "):
            blocks.append({"type": "heading1", "text": stripped[2:]})
        elif stripped.startswith("- ") or stripped.startswith("* "):
            blocks.append({"type": "bullet", "text": stripped[2:]})
        elif stripped.startswith("1. ") or (len(stripped) > 2 and stripped[0].isdigit() and stripped[1] == '.'):
            # Numbered list — find the text after "N. "
            idx = stripped.index(". ") + 2
            blocks.append({"type": "bullet", "text": stripped[idx:]})
        else:
            blocks.append({"type": "paragraph", "text": stripped})
    return blocks


def _markdown_to_docx(title: str, content: str, path: str):
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(title, level=0)

    for block in _parse_markdown_blocks(content):
        if block["type"] == "heading1":
            doc.add_heading(block["text"], level=1)
        elif block["type"] == "heading2":
            doc.add_heading(block["text"], level=2)
        elif block["type"] == "heading3":
            doc.add_heading(block["text"], level=3)
        elif block["type"] == "bullet":
            doc.add_paragraph(block["text"], style="List Bullet")
        else:
            doc.add_paragraph(block["text"])

    doc.save(path)


def _markdown_to_pptx(title: str, content: str, path: str):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title

    # Parse into slides: each # starts a new slide
    current_title = ""
    current_bullets = []

    def _flush_slide():
        if not current_title and not current_bullets:
            return
        layout = prs.slide_layouts[1]  # Title + Content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = current_title or ""
        body = s.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(current_bullets):
            if i == 0:
                tf.paragraphs[0].text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet

    for block in _parse_markdown_blocks(content):
        if block["type"] == "heading1":
            _flush_slide()
            current_title = block["text"]
            current_bullets = []
        elif block["type"] in ("heading2", "heading3"):
            current_bullets.append(block["text"])
        elif block["type"] == "bullet":
            current_bullets.append(block["text"])
        else:
            current_bullets.append(block["text"])

    _flush_slide()
    prs.save(path)


def _markdown_to_pdf(title: str, content: str, path: str):
    import fitz  # PyMuPDF

    doc = fitz.open()
    page = doc.new_page()
    y = 72  # top margin
    margin_x = 72
    max_width = page.rect.width - 2 * margin_x
    line_height = 16

    def _insert_text(text: str, fontsize: float = 11, bold: bool = False, indent: float = 0):
        nonlocal y, page
        fontname = "helv"
        # Check if we need a new page
        if y + fontsize + 8 > page.rect.height - 72:
            page = doc.new_page()
            y = 72

        # Simple word-wrap
        words = text.split()
        lines = []
        current_line = ""
        for word in words:
            test = f"{current_line} {word}".strip()
            # Approximate: ~0.5 * fontsize per char width
            if len(test) * fontsize * 0.5 > max_width - indent:
                if current_line:
                    lines.append(current_line)
                current_line = word
            else:
                current_line = test
        if current_line:
            lines.append(current_line)

        for line in lines:
            if y + fontsize + 4 > page.rect.height - 72:
                page = doc.new_page()
                y = 72
            page.insert_text(
                (margin_x + indent, y),
                line,
                fontsize=fontsize,
                fontname=fontname,
            )
            y += fontsize + 4

    # Title
    _insert_text(title, fontsize=20, bold=True)
    y += 12

    for block in _parse_markdown_blocks(content):
        if block["type"] == "heading1":
            y += 8
            _insert_text(block["text"], fontsize=16, bold=True)
            y += 4
        elif block["type"] == "heading2":
            y += 6
            _insert_text(block["text"], fontsize=14, bold=True)
            y += 2
        elif block["type"] == "heading3":
            y += 4
            _insert_text(block["text"], fontsize=12, bold=True)
        elif block["type"] == "bullet":
            _insert_text(f"•  {block['text']}", fontsize=11, indent=18)
        else:
            _insert_text(block["text"], fontsize=11)

    doc.save(path)
    doc.close()


async def _create_document(type: str, title: str, content: str) -> str:
    """Create a document file and return FILE: path for Telegram delivery."""
    _ensure_doc_dir()

    # Sanitize filename
    safe_title = "".join(c for c in title if c.isalnum() or c in " _-").strip() or "document"

    ext_map = {"word": ".docx", "slides": ".pptx", "pdf": ".pdf"}
    ext = ext_map.get(type, ".docx")
    filename = f"{safe_title}{ext}"
    path = os.path.join(_DOC_DIR, filename)

    try:
        if type == "word":
            await asyncio.get_event_loop().run_in_executor(None, _markdown_to_docx, title, content, path)
        elif type == "slides":
            await asyncio.get_event_loop().run_in_executor(None, _markdown_to_pptx, title, content, path)
        elif type == "pdf":
            await asyncio.get_event_loop().run_in_executor(None, _markdown_to_pdf, title, content, path)
        else:
            return f"不支持的文档类型: {type}"
    except Exception as e:
        return f"创建文档失败: {e}"

    return f"FILE:{path}\n已创建文档: {filename}"


_TOOL_MAP = {
    "read_emails": _read_emails,
    "search_emails": _search_emails,
    "send_email": _send_email,
    "capture_screen": _capture_screen,
    "check_wechat": _check_wechat,
    "run_claude_code": _run_claude_code,
    "find_files": _find_files,
    "read_file": _read_file,
    "open_app": _open_app,
    "quit_app": _quit_app,
    "run_command": _run_command,
    "get_calendar_events": _get_calendar_events,
    "create_reminder": _create_reminder,
    "music_play": _music_play,
    "music_pause": _music_pause,
    "music_next": _music_next,
    "music_previous": _music_previous,
    "music_now_playing": _music_now_playing,
    "music_search_play": _music_search_play,
    "create_document": _create_document,
}
